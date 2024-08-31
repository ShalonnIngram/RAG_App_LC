import io
import os
import docx
import csv
import pandas as pd
from PyPDF2 import PdfReader
from pptx import Presentation
import boto3
from dotenv import load_dotenv
import streamlit as st
import openai
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import Chroma
from langchain import OpenAI, VectorDBQA
from langchain.chains import RetrievalQAWithSourcesChain

# Load environment variables
load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")


# Utility function to read and extract text from various document types
def read_and_textify(files):
    text_list = []
    sources_list = []
    
    for file in files:
        file_extension = file.name.split('.')[-1].lower()
        
        try:
            file_buffer = io.BytesIO(file.read())
            
            if file_extension == "pdf":
                pdf_reader = PdfReader(file_buffer)
                for i, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    text_list.append(text)
                    sources_list.append(f"{file.name}_page_{i}")
            
            elif file_extension == "docx":
                doc = docx.Document(file_buffer)
                text = "\n".join([para.text for para in doc.paragraphs])
                text_list.append(text)
                sources_list.append(file.name)
                
            elif file_extension == "pptx":
                prs = Presentation(file_buffer)
                text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
                text_list.append(text)
                sources_list.append(file.name)
            
            elif file_extension == "txt":
                text = file_buffer.read().decode("utf-8")
                text_list.append(text)
                sources_list.append(file.name)
            
            elif file_extension == "csv":
                reader = csv.reader(io.StringIO(file_buffer.read().decode("utf-8")))
                text = "\n".join([", ".join(row) for row in reader])
                text_list.append(text)
                sources_list.append(file.name)
            
            elif file_extension == "xlsx":
                df = pd.read_excel(file_buffer)
                text = df.to_csv(index=False)
                text_list.append(text)
                sources_list.append(file.name)
            
            # Add additional file types as needed

        except Exception as e:
            print(f"Error processing file {file.name}: {e}")
    
    return text_list, sources_list

# Function to initialize Chroma vectorstore
def initialize_vectorstore(documents, sources):
    embeddings = OpenAIEmbeddings()
    vstore = Chroma.from_texts(documents, embeddings, metadatas=[{"source": s} for s in sources])
    return vstore

# Streamlit application
def main():
    st.set_page_config(page_title="LixCap Internal Document Query App")
    st.header("LixCap Internal Document Query Application")
    st.write("---")
    
    with st.sidebar:
        st.image("logo.png")
        st.subheader("For Internal Use Only")

    uploaded_files = st.file_uploader("Please upload files", accept_multiple_files=True, type=["txt", "pdf", "docx", "pptx", "csv", "xlsx"])
    st.write("---")

    if uploaded_files:
        st.write(f"{len(uploaded_files)} document(s) loaded..")
        textify_output = read_and_textify(uploaded_files)
        documents, sources = textify_output

        if documents:
            vstore = initialize_vectorstore(documents, sources)
            
            model_name = "gpt-4o-mini"
            retriever = vstore.as_retriever()
            retriever.search_kwargs = {'k': 2}

            llm = OpenAI(model_name=model_name, streaming=True)
            model = RetrievalQAWithSourcesChain.from_chain_type(llm=llm, chain_type="stuff", retriever=retriever)
            
            st.header("Ask Questions about your Documents")
            user_q = st.text_area("Please submit your questions below")

            if st.button("Get Response"):
                if user_q:
                    try:
                        with st.spinner("Processing..."):
                            result = model({"question": user_q}, return_only_outputs=True)
                            st.subheader('Your response:')
                            st.write(result['answer'])
                            st.subheader('Source pages:')
                            st.write(result['sources'])
                    except Exception as e:
                        st.error(f"An error occurred: {e}")
                        st.error('Oops, the GPT response resulted in an error. Please try again with a different question.')
                else:
                    st.warning("Please enter a question to get a response.")
        else:
            st.warning("No text found in the uploaded documents.")
    else:
        st.info("Upload files to analyze.")

if __name__ == "__main__":
    main()
